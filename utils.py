import os
import io
import json
import re
import requests
from pptx.util import Pt
import matplotlib.pyplot as plt
from pptx import Presentation
from docx import Document
from dotenv import load_dotenv
import math
from io import BytesIO
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
import uuid

load_dotenv()

# ----------------------------
# 환경 변수 (Azure 설정)
# ----------------------------
AZURE_OPENAI_ENDPOINT = os.getenv("AZURE_OPENAI_ENDPOINT")
AZURE_OPENAI_API_KEY = os.getenv("AZURE_OPENAI_API_KEY")
AZURE_OPENAI_DEPLOYMENT = os.getenv("AZURE_OPENAI_DEPLOYMENT")

AZURE_SEARCH_ENDPOINT = os.getenv("AZURE_SEARCH_ENDPOINT")
AZURE_SEARCH_KEY = os.getenv("AZURE_SEARCH_KEY")
AZURE_SEARCH_INDEX = os.getenv("AZURE_SEARCH_INDEX")
AZURE_SEARCH_API_VERSION = os.getenv("AZURE_SEARCH_API_VERSION")
AZURE_SEARCH_SEMANTIC_NAME = os.getenv("AZURE_SEARCH_SEMANTIC_NAME")  # 있으면 세맨틱

# ----------------------------
# OpenAI 호출
# ----------------------------
def _assert_aoai_ready():
    miss = []
    if not AZURE_OPENAI_ENDPOINT or not AZURE_OPENAI_ENDPOINT.startswith("https://"):
        miss.append("AZURE_OPENAI_ENDPOINT(https://...)")
    if not AZURE_OPENAI_API_KEY:
        miss.append("AZURE_OPENAI_API_KEY")
    if not AZURE_OPENAI_DEPLOYMENT:
        miss.append("AZURE_OPENAI_DEPLOYMENT(배포명)")
    if miss:
        raise RuntimeError("환경변수 누락: " + ", ".join(miss))


def call_openai(messages, max_tokens=800, temperature=0.2, force_json=True):
    _assert_aoai_ready()
    url = f"{AZURE_OPENAI_ENDPOINT}/openai/deployments/{AZURE_OPENAI_DEPLOYMENT}/chat/completions?api-version=2024-02-15-preview"
    headers = {"Content-Type": "application/json", "api-key": AZURE_OPENAI_API_KEY}
    body = {"messages": messages, "temperature": temperature, "max_tokens": max_tokens}
    if force_json:
        body["response_format"] = {"type": "json_object"}
    r = requests.post(url, headers=headers, json=body, timeout=60)
    if r.status_code != 200:
        raise RuntimeError(f"AOAI 호출 실패 {r.status_code}: {r.text[:500]}")
    data = r.json()
    return data["choices"][0]["message"]["content"]

def ready_openai() -> bool:
    return (
        AZURE_OPENAI_ENDPOINT.startswith("https://")
        and bool(AZURE_OPENAI_API_KEY)
        and bool(AZURE_OPENAI_DEPLOYMENT)
    )

def ready_search() -> bool:
    return (
        AZURE_SEARCH_ENDPOINT.startswith("https://")
        and bool(AZURE_SEARCH_KEY)
        and bool(AZURE_SEARCH_INDEX)
    )

# ----------------------------
# Azure Cognitive Search 호출
# ----------------------------
# ---------------------------
# Azure AI Search
# ---------------------------
def search_topk(query: str, k: int = 3):
    if not ready_search():
        return []

    base = f"{AZURE_SEARCH_ENDPOINT}/indexes/{AZURE_SEARCH_INDEX}/docs/search"
    params = {"api-version": AZURE_SEARCH_API_VERSION}
    headers = {"Content-Type": "application/json", "api-key": AZURE_SEARCH_KEY}
    q = query.strip() if query and query.strip() else "*"

    def _simple():
        payload = {"search": q, "top": k}
        r = requests.post(base, headers=headers, params=params, json=payload, timeout=30)
        r.raise_for_status()
        return r

    # 시맨틱 설정이 있을 때만 시도
    if AZURE_SEARCH_SEMANTIC_NAME:
        payload_sem = {
            "search": q,
            "top": k,
            "queryType": "semantic",
            "queryLanguage": "ko",
            "semanticConfiguration": AZURE_SEARCH_SEMANTIC_NAME,
        }
        try:
            r = requests.post(base, headers=headers, params=params, json=payload_sem, timeout=30)
            if r.status_code != 200:
                r = _simple()
        except Exception:
            r = _simple()
    else:
        r = _simple()

    items = r.json().get("value", [])
    return [{
        "id": it.get("id"),
        "title": it.get("title", ""),
        "category": it.get("category", ""),
        "created_at": it.get("created_at", ""),
        "content": it.get("content", "")[:1200],
    } for it in items]\
    
def aoai_chat(messages, max_tokens=900, temperature=0.6) -> str:
    """
    일반 텍스트 응답. 실패 시 HTTPError 발생.
    """
    if not ready_openai():
        raise RuntimeError("(OpenAI 환경변수 미설정: .env 확인)")

    url = f"{AZURE_OPENAI_ENDPOINT}/openai/deployments/{AZURE_OPENAI_DEPLOYMENT}/chat/completions"
    params = {"api-version": "2024-02-15-preview"}
    headers = {"api-key": AZURE_OPENAI_API_KEY, "Content-Type": "application/json"}
    payload = {"messages": messages, "max_tokens": max_tokens, "temperature": temperature}

    r = requests.post(url, headers=headers, params=params, json=payload, timeout=60)
    r.raise_for_status()
    data = r.json()
    return data["choices"][0]["message"]["content"]

# ----------------------------
# 개요 생성 + PPT 만들기
# ----------------------------
def generate_outline_and_ppt(
    topic: str,
    keypoints: str,
    audience: str,
    tone: str,
    slide_count: int,
    use_rag: bool,
    template_file,              # st.file_uploader 결과 또는 None
    deck_title: str,
    font_name: str = "맑은 고딕"
):
    """
    1) LLM으로 개요(JSON 유도) 생성
    2) 안전 정규화 → 표준 구조로 변환
    3) 템플릿을 적용해 PPT 생성
    반환: (ppt_bytes: BytesIO, outline_list: list[dict])
    """
    # --- RAG 컨텍스트(선택)
    context = ""
    try:
        if use_rag:
            # 프로젝트에서 쓰던 검색 함수가 있으면 활용 (없어도 진행)
            from utils import search_references  # 있으면
            refs = search_references(topic, top_k=3) or []
            context = "\n".join([r.get("snippet", "") for r in refs])
    except Exception:
        pass

    # --- LLM 호출
    sys = "You are a Korean PPT planner. Output JSON only."
    usr = f"""
다음 정보를 바탕으로 PPT 개요를 생성하라.
각 슬라이드는 다음 JSON 원소 형식만 사용:
{{"title":"...", "bullets":["...", "..."]}}

제목/핵심내용/대상/톤/페이지수를 반영할 것.
필요하면 '결론/요약' 슬라이드 포함.

주제: {topic}
핵심내용: {keypoints}
대상: {audience}
톤: {tone}
슬라이드 수: {slide_count}

참고 문서 발췌(선택):
{context or "(없음)"}
"""
    try:
        raw = call_openai(
            [{"role": "system", "content": sys}, {"role": "user", "content": usr}],
            max_tokens=1200, temperature=0.2, force_json=True
        )
    except Exception as e:
        # LLM 실패해도 빈 개요로 진행
        raw = ""

    # --- 정규화
    try:
        # 우선 안전 JSON 파싱 시도
        parsed = _parse_json_safely(raw)
    except Exception:
        parsed = raw
    outline = _normalize_outline(parsed, slide_count)

    # --- PPT 생성
    try:
        prs = Presentation(template_file) if template_file else Presentation()
    except Exception:
        prs = Presentation()

    # 안전한 레이아웃 선택
    title_layout = None
    content_layout = None
    if prs.slide_layouts:
        # 0: Title, 1: Title and Content (보통)
        title_layout = prs.slide_layouts[0] if len(prs.slide_layouts) > 0 else None
        content_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else None

    # 타이틀 슬라이드
    if title_layout:
        s0 = prs.slides.add_slide(title_layout)
        try:
            if s0.shapes.title:
                s0.shapes.title.text = deck_title or topic or "프로젝트 제안"
            if len(s0.placeholders) > 1 and hasattr(s0.placeholders[1], "text_frame"):
                s0.placeholders[1].text = audience or ""
        except Exception:
            pass

    # 본문 슬라이드들
    for sl in outline:
        # 혹시 모를 문자열 방지 (핵심 버그 방지)
        if isinstance(sl, str):
            sl = {"title": sl, "bullets": []}
        title = sl.get("title") or "제목 없음"
        bullets = sl.get("bullets") or []
        if content_layout:
            slide = prs.slides.add_slide(content_layout)
        else:
            slide = prs.slides.add_slide(prs.slide_layouts[0])
        # 제목
        try:
            if slide.shapes.title:
                slide.shapes.title.text = str(title)
        except Exception:
            pass
        # 본문
        try:
            # Title and Content 레이아웃의 본문 플레이스홀더 찾기
            body = None
            for ph in slide.placeholders:
                if hasattr(ph, "text_frame") and ph != slide.shapes.title:
                    body = ph
                    break
            if body is not None:
                tf = body.text_frame
                tf.clear()
                if bullets:
                    tf.text = str(bullets[0])
                    for b in bullets[1:]:
                        p = tf.add_paragraph()
                        p.text = str(b)
                        p.level = 0
                else:
                    tf.text = ""
        except Exception:
            pass

    # 폰트 통일(간단 적용)
    try:
        if font_name:
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame"):
                        for p in shape.text_frame.paragraphs:
                            for r in p.runs:
                                r.font.name = font_name
                                r.font.size = Pt(18)
    except Exception:
        pass

    bio = io.BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio, outline

# ---------------------------
# 개요 텍스트 파서 (마크다운/일반 텍스트)
# ---------------------------
def parse_outline_from_text(text: str):
    """
    마크다운/일반 텍스트 개요 → [{"title":..., "bullets":[...]}] 리스트로 변환
    지원 예:
      ## 슬라이드 1: 제목
      - 포인트
      - 포인트
    """
    slides = []
    lines = [ln.rstrip() for ln in text.splitlines()]

    title_pat = re.compile(
        r'^(?:#{1,3}\s*|(?:\d+[\.\)]\s*)|(?:슬라이드\s*\d+\s*[:\-]\s*)|(?:slide\s*\d+\s*[:\-]\s*)|^)(?:슬라이드\s*\d+\s*[:\-]\s*)?(.*)$',
        re.IGNORECASE
    )

    current = {"title": None, "bullets": []}

    def push_current():
        if current["title"] and (current["bullets"] or current["title"].strip()):
            slides.append({
                "title": current["title"].strip(),
                "bullets": [b.strip() for b in current["bullets"] if b.strip()]
            })

    for ln in lines:
        if not ln.strip():
            continue

        m = title_pat.match(ln)
        is_title_line = False
        if m:
            raw = m.group(1).strip()
            if not raw.startswith(("-", "*", "•")) and len(raw) > 0:
                if current["title"] is not None:
                    push_current()
                    current = {"title": None, "bullets": []}
                current["title"] = raw
                is_title_line = True

        if is_title_line:
            continue

        if ln.lstrip().startswith(("-", "*", "•", "·")):
            bullet = ln.lstrip()[1:].strip()
            if current["title"] is None:
                current["title"] = "슬라이드"
            current["bullets"].append(bullet)
        else:
            if current["title"] is None:
                current["title"] = ln.strip()
            else:
                current["bullets"].append(ln.strip())

    push_current()

    if not slides and text.strip():
        slides = [{"title": "개요", "bullets": [text.strip()[:1800]]}]

    return slides

def _apply_font_safe(shape_or_paragraph, font_name, size):
    """안전한 폰트 적용"""
    try:
        if hasattr(shape_or_paragraph, 'text_frame'):
            # Shape인 경우
            for paragraph in shape_or_paragraph.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = font_name
                    run.font.size = Pt(size)
        elif hasattr(shape_or_paragraph, 'runs'):
            # Paragraph인 경우
            for run in shape_or_paragraph.runs:
                run.font.name = font_name
                run.font.size = Pt(size)
        elif hasattr(shape_or_paragraph, 'font'):
            # Run인 경우
            shape_or_paragraph.font.name = font_name
            shape_or_paragraph.font.size = Pt(size)
    except Exception as e:
        print(f"폰트 적용 오류: {e}")

def _add_title_only_slide_safe(prs, title_text, layout_info, font="Malgun Gothic"):
    """안전한 타이틀 전용 슬라이드 추가"""
    try:
        layout_idx = layout_info.get('section') or layout_info.get('title_only') or layout_info['title']
        layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(layout)
        print(f"섹션 슬라이드 생성: 레이아웃 {layout_idx} 사용")
        
        # 타이틀 placeholder 찾기
        for shape in slide.shapes:
            if shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                if ph_type in [1, 3]:  # TITLE, CENTER_TITLE
                    shape.text = title_text
                    _apply_font_to_shape(shape, font, 36)
                    break
        else:
            # placeholder가 없다면 첫 번째 텍스트 shape 사용
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame'):
                    shape.text = title_text
                    _apply_font_to_shape(shape, font, 36)
                    break
        
        return slide
    except Exception as e:
        print(f"섹션 슬라이드 생성 오류: {e}")
        return None

def _apply_font_to_shape(shape, font_name, size):
    """Shape의 모든 텍스트에 폰트 적용"""
    try:
        if hasattr(shape, 'text_frame'):
            for paragraph in shape.text_frame.paragraphs:
                _apply_font_to_paragraph(paragraph, font_name, size)
    except Exception as e:
        print(f"Shape 폰트 적용 오류: {e}")

def _apply_font_to_paragraph(paragraph, font_name, size):
    """Paragraph의 모든 run에 폰트 적용"""
    try:
        # 기존 run이 없으면 생성
        if not paragraph.runs:
            paragraph._p.get_or_add_r()
        
        for run in paragraph.runs:
            if hasattr(run, 'font'):
                run.font.name = font_name
                run.font.size = Pt(size)
    except Exception as e:
        print(f"Paragraph 폰트 적용 오류: {e}")
    
def _add_content_slide_safe(prs, title, bullets, layout_info, font="Malgun Gothic"):
    """템플릿의 콘텐츠 레이아웃으로 콘텐츠 슬라이드 생성"""
    try:
        layout = prs.slide_layouts[layout_info['content']]
        slide = prs.slides.add_slide(layout)
        print(f"콘텐츠 슬라이드 생성: 레이아웃 {layout_info['content']} 사용")
        
        # placeholder 분류
        title_placeholder = None
        content_placeholder = None
        
        for shape in slide.shapes:
            if shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                if ph_type in [1, 3]:  # TITLE, CENTER_TITLE
                    title_placeholder = shape
                elif ph_type in [2, 7, 8, 13]:  # BODY, FOOTER, HEADER, CONTENT
                    if content_placeholder is None:  # 첫 번째 콘텐츠 placeholder 사용
                        content_placeholder = shape
        
        # 타이틀 설정
        if title_placeholder:
            title_placeholder.text = title
            _apply_font_to_shape(title_placeholder, font, 32)
        
        # 콘텐츠 설정
        if content_placeholder and hasattr(content_placeholder, 'text_frame'):
            tf = content_placeholder.text_frame
            tf.clear()
            
            if bullets:
                # 첫 번째 불릿
                if tf.paragraphs:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = str(bullets[0])
                p.level = 0
                _apply_font_to_paragraph(p, font, 18)
                
                # 나머지 불릿들
                for bullet in bullets[1:]:
                    p = tf.add_paragraph()
                    p.text = str(bullet)
                    p.level = 0
                    _apply_font_to_paragraph(p, font, 18)
        
        return slide
    except Exception as e:
        print(f"콘텐츠 슬라이드 생성 오류: {e}")
        return None

def _analyze_layouts(prs):
    """템플릿의 레이아웃을 분석하여 최적의 레이아웃을 찾음"""
    layouts = {
        'title': None,
        'title_only': None,
        'content': None,
        'section': None
    }
    
    for i, layout in enumerate(prs.slide_layouts):
        name_lower = layout.name.lower()
        print(f"레이아웃 {i}: {layout.name}")
        
        # 타이틀 슬라이드 찾기
        if any(keyword in name_lower for keyword in ['title slide', '제목 슬라이드', 'title']):
            if 'only' not in name_lower and layouts['title'] is None:
                layouts['title'] = i
        
        # 타이틀 온리 찾기
        if any(keyword in name_lower for keyword in ['title only', '제목만', 'section']):
            layouts['title_only'] = i
        
        # 콘텐츠 슬라이드 찾기
        if any(keyword in name_lower for keyword in ['content', 'bullet', '내용', 'two content']):
            layouts['content'] = i
        
        # 섹션 헤더 찾기
        if any(keyword in name_lower for keyword in ['section', '섹션', 'divider']):
            layouts['section'] = i
    
    # 기본값 설정
    if layouts['title'] is None:
        layouts['title'] = 0  # 첫 번째 레이아웃
    if layouts['content'] is None:
        layouts['content'] = min(1, len(prs.slide_layouts) - 1)  # 두 번째 레이아웃
    if layouts['title_only'] is None:
        layouts['title_only'] = layouts['title']  # 타이틀과 동일
    if layouts['section'] is None:
        layouts['section'] = layouts['title_only']  # 타이틀 온리와 동일
    
    return layouts
    
def build_ppt_from_outline_mixed(
    outline_slides, project_title=None, template_bytes=None, font_name="Malgun Gothic"
) -> BytesIO:
    """
    템플릿의 디자인과 배경을 완전히 유지하면서 콘텐츠만 변경하는 PPT 생성
    - 첫 장: 타이틀
    - 섹션: bullets 없음 or '섹션:' 등으로 시작
    - 나머지: 제목+불릿
    """
    try:
        # 템플릿 처리 개선
        if template_bytes:
            if hasattr(template_bytes, 'read'):
                # file_uploader 객체인 경우
                template_data = template_bytes.read()
                template_bytes.seek(0)  # 포인터 리셋
            else:
                # 이미 bytes인 경우
                template_data = template_bytes
            
            prs = Presentation(BytesIO(template_data))

            # 기존 슬라이드가 있다면 모두 삭제 (템플릿 구조는 유지)
            slide_ids = [slide.slide_id for slide in prs.slides]
            for slide_id in slide_ids:
                for i, slide in enumerate(prs.slides):
                    if slide.slide_id == slide_id:
                        xml_slides = prs.slides._sldIdLst
                        xml_slides.remove(xml_slides[i])
                        break

        else:
            prs = Presentation()
            print("기본 템플릿 사용")
    except Exception as e:
        print(f"템플릿 로드 실패, 기본 템플릿 사용: {e}")
        prs = Presentation()

    # 사용 가능한 레이아웃 분석
    layout_info = _analyze_layouts(prs)
    print(f"레이아웃 분석 결과: {layout_info}")

    # 1) 타이틀 슬라이드
    title_text = project_title or "Project Pitch"
    _add_title_only_slide(prs, title_text, layout_info, font=font_name)

    # 2) 본문 슬라이드들
    for item in outline_slides:
        if not isinstance(item, dict):
            continue
            
        title = item.get("title") or "제목 없음"
        bullets = item.get("bullets") or []
        
        if _is_section_slide(item):
            _add_title_only_slide_safe(prs, title, layout_info, font=font_name)
        else:
            _add_content_slide_safe(prs, title, bullets, layout_info, font=font_name)

    bio = BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio


def _is_section_slide(item) -> bool:
    """섹션 슬라이드 판별 함수"""
    title = (item.get("title") or "").strip()
    bullets = item.get("bullets") or []
    if not bullets:
        return True
    lowered = title.lower()
    return (
        lowered.startswith("섹션:")
        or lowered.startswith("[섹션]")
        or lowered.startswith("section:")
        or lowered.startswith("section ")
        or lowered.startswith("section-")
        or "목차" in lowered
        or "overview" in lowered
    )

def _add_title_only_slide(prs, title_text, layout_info, font="Malgun Gothic"):
    """템플릿의 타이틀 슬라이드 레이아웃으로 타이틀 슬라이드 생성"""
    try:
        layout_idx = layout_info.get('section') or layout_info.get('title_only') or layout_info['title']
        layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(layout)
        print(f"섹션 슬라이드 생성: 레이아웃 {layout_idx} 사용")
        
        # 타이틀 placeholder 찾기
        for shape in slide.shapes:
            if shape.is_placeholder:
                ph_type = shape.placeholder_format.type
                if ph_type in [1, 3]:  # TITLE, CENTER_TITLE
                    shape.text = title_text
                    _apply_font_to_shape(shape, font, 36)
                    break
        else:
            # placeholder가 없다면 첫 번째 텍스트 shape 사용
            for shape in slide.shapes:
                if hasattr(shape, 'text_frame'):
                    shape.text = title_text
                    _apply_font_to_shape(shape, font, 36)
                    break
        
        return slide
    except Exception as e:
        print(f"섹션 슬라이드 생성 오류: {e}")
        return None

# ---------------------------
# PPT 생성 (타이틀/섹션/콘텐츠 혼합)
# ---------------------------
def _find_title_body_placeholders(slide):
    title_ph, body_ph = None, None
    for shape in slide.shapes:
        if getattr(shape, "is_placeholder", False):
            phf = shape.placeholder_format
            if phf and phf.type in [1, 3]:  # TITLE, CENTER_TITLE
                title_ph = shape
            if phf and phf.type in [2, 4, 7, 8, 9, 10]:  # BODY, SUBTITLE 등
                if body_ph is None:
                    body_ph = shape

    if title_ph is None:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                title_ph = shape
                break
    if body_ph is None:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape is not title_ph:
                body_ph = shape
                break
    return title_ph, body_ph

def _set_text(shape, text, font="Malgun Gothic", size=28, align=PP_PARAGRAPH_ALIGNMENT.LEFT):
    if shape is None or not getattr(shape, "has_text_frame", False):
        return
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text or ""
    p.font.name = font
    p.font.size = Pt(size)
    p.alignment = align

def _normalize_outline(raw, slide_count: int):
    """
    raw(LLM 응답)가 문자열/문자열 리스트/딕트 리스트 어떤 형태든
    [{ "title": str, "bullets": [str, ...] }, ...] 로 변환.
    """
    # 1) 문자열이면 JSON 파싱 시도 → 실패하면 라인 파싱
    if isinstance(raw, str):
        try:
            raw = json.loads(raw)
        except Exception:
            # 라인 파서: 숫자.제목 / - bullet
            lines = [l.strip() for l in raw.splitlines() if l.strip()]
            slides = []
            cur = None
            for ln in lines:
                if ln[:2].isdigit() or ln[0].isdigit() or ln.startswith("슬라이드"):
                    # 새 제목 추정
                    if cur:
                        slides.append(cur)
                    # 숫자/접두 제거
                    title = ln
                    title = title.split(" ", 1)[-1] if ". " in title else title.replace("슬라이드", "").strip(": ")
                    cur = {"title": title.strip(), "bullets": []}
                elif ln.startswith("-") or ln.startswith("•") or ln.startswith("*"):
                    if not cur:
                        cur = {"title": "제목 없음", "bullets": []}
                    cur["bullets"].append(ln.lstrip("-•* ").strip())
                else:
                    # 제목 라인으로 취급
                    if cur:
                        slides.append(cur)
                    cur = {"title": ln.strip(), "bullets": []}
            if cur:
                slides.append(cur)
            raw = slides

    # 2) 리스트이면 각 요소 정규화
    if isinstance(raw, list):
        norm = []
        for item in raw:
            if isinstance(item, str):
                norm.append({"title": item.strip(), "bullets": []})
            elif isinstance(item, dict):
                title = str(item.get("title") or item.get("heading") or "제목 없음").strip()
                bullets = item.get("bullets") or item.get("points") or []
                if isinstance(bullets, str):
                    bullets = [bullets]
                # 불릿들은 문자열만 남기기
                bullets = [str(b).strip() for b in bullets if str(b).strip()]
                norm.append({"title": title, "bullets": bullets})
            else:
                # 알 수 없는 타입 → 문자열화
                norm.append({"title": str(item), "bullets": []})
        # 슬라이드 수 제한
        return norm[:max(1, slide_count)]
    else:
        # 완전 예외 → 기본 한 장
        return [{"title": "개요", "bullets": []}]

# ----------------------------
# 발표 시뮬레이션 라운드
# ----------------------------
def run_simulation_round(user_input, persona, style, intensity):
    messages = [
        {"role": "system", "content": f"너는 {persona} 역할의 심사위원이다. 성향: {style}, 강도 {intensity}"},
        {"role": "user", "content": f"발표자가 이렇게 답변함: {user_input}. 이에 대한 질문을 해줘."}
    ]
    reply = call_openai(messages)

    eval_messages = [
        {"role": "system", "content": "너는 발표 심사위원이다."},
        {"role": "user", "content": f"발표자 답변: {user_input}\n심사위원 질문: {reply}\n답변을 평가해라. JSON으로 {{설득력, 전문성, 소통력, 압박대응, 총평}} 반환"}
    ]
    eval_result = json.loads(call_openai(eval_messages))

    return reply, eval_result

def _assert_aoai_ready():
    miss = []
    if not AZURE_OPENAI_ENDPOINT or not AZURE_OPENAI_ENDPOINT.startswith("https://"):
        miss.append("AZURE_OPENAI_ENDPOINT(https://...)")
    if not AZURE_OPENAI_API_KEY:
        miss.append("AZURE_OPENAI_API_KEY")
    if not AZURE_OPENAI_DEPLOYMENT:
        miss.append("AZURE_OPENAI_DEPLOYMENT(배포명)")
    if miss:
        raise RuntimeError("환경변수 누락: " + ", ".join(miss))

def _parse_json_safely(text: str):
    """
    LLM 응답에서 JSON만 뽑아 파싱.
    1) 그대로 로드
    2) ```json ... ``` 블록 제거
    3) 첫 { ~ 마지막 } 추출 (정규식 포함)
    실패 시 상세 메시지와 함께 예외
    """
    import json, re

    if text is None:
        raise ValueError("LLM 응답이 비었습니다(None).")
    t = text.strip()

    # 1) 바로 시도
    try:
        return json.loads(t)
    except Exception:
        pass

    # 2) 코드펜스 제거
    m = re.search(r"```(?:json)?\s*(.+?)```", t, flags=re.S)
    if m:
        inner = m.group(1).strip()
        try:
            return json.loads(inner)
        except Exception:
            # inner에서 중괄호만 재시도
            start = inner.find("{")
            end = inner.rfind("}")
            if start != -1 and end != -1 and end > start:
                core = inner[start:end+1]
                try:
                    return json.loads(core)
                except Exception:
                    pass

    # 3) 가장 바깥 { ... } 추출 (정규식 보강)
    match = re.search(r"\{.*\}", t, re.S)
    if match:
        core = match.group(0)
        try:
            return json.loads(core)
        except Exception:
            pass

    # 디버깅용 원문 일부 남기고 실패
    snippet = t[:400].replace("\n", " ")
    raise ValueError(f"JSON 파싱 실패. 응답 앞부분: {snippet}")


def call_openai(messages, max_tokens=800, temperature=0.2, force_json=True):
    """
    Azure OpenAI Chat Completions 호출.
    - force_json=True면 JSON-only 유도(response_format 지원 / 프롬프트 강제)
    - 실패 시 상세 에러 텍스트 포함해 예외
    """
    _assert_aoai_ready()

    url = f"{AZURE_OPENAI_ENDPOINT}/openai/deployments/{AZURE_OPENAI_DEPLOYMENT}/chat/completions?api-version=2024-02-15-preview"
    headers = {"Content-Type": "application/json", "api-key": AZURE_OPENAI_API_KEY}
    body = {
        "messages": messages,
        "temperature": temperature,
        "max_tokens": max_tokens,
    }
    # 일부 배포는 response_format JSON 지원. 미지원이어도 프롬프트에서 JSON만 요구하니 문제 없음.
    if force_json:
        body["response_format"] = {"type": "json_object"}

    r = requests.post(url, headers=headers, json=body, timeout=60)
    if r.status_code != 200:
        # 에러 바디 그대로 노출(디버깅 편의)
        raise RuntimeError(f"AOAI 호출 실패 {r.status_code}: {r.text[:500]}")

    data = r.json()
    try:
        content = data["choices"][0]["message"]["content"]
    except Exception:
        raise RuntimeError(f"AOAI 응답 형식 예기치 않음: {json.dumps(data)[:500]}")

    return content

def search_references(query: str, top_k: int = 3):
    """
    Azure Cognitive Search에서 관련 문헌 상위 top_k 조각을 가져온다.
    반환: [{"source": str, "snippet": str}, ...]
    """
    if not (AZURE_SEARCH_ENDPOINT and AZURE_SEARCH_KEY and AZURE_SEARCH_INDEX):
        return []

    url = f"{AZURE_SEARCH_ENDPOINT}/indexes/{AZURE_SEARCH_INDEX}/docs/search?api-version={AZURE_SEARCH_API_VERSION}"
    headers = {"Content-Type": "application/json", "api-key": AZURE_SEARCH_KEY}

    body = {
        "search": query if query else "*",
        "top": top_k
    }
    # 세맨틱 구성 있으면 살짝 도움 (없어도 동작)
    if AZURE_SEARCH_SEMANTIC_NAME:
        body["queryType"] = "semantic"
        body["semanticConfiguration"] = AZURE_SEARCH_SEMANTIC_NAME

    try:
        r = requests.post(url, headers=headers, json=body, timeout=15)
        r.raise_for_status()
        js = r.json()
        out = []
        for i, v in enumerate(js.get("value", []), start=1):
            name = _pick_display_source(v)
            if not name:  # 난수/ID만 있을 때
                name = f"문서 {i:03d}"
            out.append({
                "display": name,
                "url": _pick_url(v),
                "snippet": (_pick_snippet(v) or "").replace("\n", " "),
                "raw": v,
            })
        return out
    except Exception as e:
        print("RAG search error:", e)
        return []
    
def build_rag_for_slides(slides, top_k=2):
    """
    각 슬라이드(title+content)를 쿼리로 삼아 RAG refs를 붙인다.
    반환: { idx: [{"source":..,"snippet":..}, ...], ... }
    """
    rag_map = {}
    for s in slides:
        q = (s.get("title","") + " " + s.get("content","")).strip()
        refs = search_references(q, top_k=top_k) if q else []
        rag_map[s["idx"]] = refs
    return rag_map

# ----------------------------
# PPT 업로드 평가
# ----------------------------
def analyze_ppt_with_llm(uploaded_ppt, criteria, use_rag, rag_top_k):
    """
    LLM Judge + (옵션) RAG 근거 기반 평가.
    반환: (scores_dict, feedback_top_list, rich_result_dict)
    """
    import json
    from pptx import Presentation

    # ---- 0) PPT 텍스트 추출 (중복 append 제거) ----
    prs = Presentation(uploaded_ppt)
    slides = []
    for idx, slide in enumerate(prs.slides, start=1):
        try:
            s = _extract_slide(slide, idx)  # 그룹/테이블 포함 추출 버전 권장
            if not isinstance(s, dict) or "idx" not in s:
                raise ValueError("extract returned invalid")
            slides.append(s)
        except Exception:
            title = ""
            if getattr(slide.shapes, "title", None):
                try:
                    title = slide.shapes.title.text or ""
                except Exception:
                    title = ""
            body_chunks = []
            for shp in slide.shapes:
                try:
                    if shp.has_text_frame and shp != getattr(slide.shapes, "title", None):
                        body_chunks.append(shp.text.strip())
                except Exception:
                    continue
            slides.append({
                "idx": idx,
                "title": (title or f"슬라이드 {idx}").strip(),
                "content": " ".join([c for c in body_chunks if c]).strip()
            })

    # ---- 1) 입력 길이 제한 (LLM 안정성) ----
    MAX_CHARS = 6000
    total = 0
    trimmed = []
    for s in slides:
        t = s.get("content", "") or ""
        if total + len(t) > MAX_CHARS:
            remain = max(0, MAX_CHARS - total)
            t = t[:remain]
        total += len(t)
        trimmed.append({"idx": s["idx"], "title": s.get("title",""), "content": t})
        if total >= MAX_CHARS:
            break
    slides_for_llm = trimmed if trimmed else slides

    # ---- 2) 기준/가중치/루브릭 문자열 ----
    crit_lines, weight_map = [], {}
    for c in criteria:
        name = (str(c.get("name","")) or "기준").strip()
        weight = int(c.get("weight", 0))
        rubric = (str(c.get("rubric","")) or "").strip()
        crit_lines.append(f"- {name}({weight}%): {rubric}")
        weight_map[name] = max(0, weight)
    crit_block = "\n".join(crit_lines) if crit_lines else "- 전반 평가(100%)"
    denom = sum(weight_map.values()) or 1

    # ---- 3) (옵션) RAG 참조 생성 ----
    rag_map = build_rag_for_slides(slides_for_llm, top_k=rag_top_k) if use_rag else {}
    rag_text = ""
    if use_rag and rag_map:
        rag_lines = []
        for s in slides_for_llm:
            sid = s["idx"]
            refs = rag_map.get(sid, [])
            if not refs:
                continue
            rag_lines.append(f"[Slide {sid}: {s.get('title','')}]")
            for i, r in enumerate(refs, start=1):
                src = r.get("source","")
                snip = (r.get("snippet","") or "").replace("\n", " ")
                rag_lines.append(f"  - ({i}) {src}: {snip[:300]}")
        rag_text = "\n".join(rag_lines)

    # ---- 4) 평가 프롬프트 (RAG 포함) ----
    sys = "You are a rigorous Korean PPT evaluator. Output JSON only."
    usr = f"""
아래 PPT 슬라이드 요약을 '사용자 정의 루브릭'에 따라 평가하라.
가능하면 제공된 RAG 참고 문헌과 비교/검증하라.
한국어로, 설명 없이 JSON만 출력한다.

[슬라이드 요약(JSON)]
{json.dumps(slides_for_llm, ensure_ascii=False)}

[평가기준(가중치/루브릭)]
{crit_block}

[참고 문헌 발췌(RAG)]
{rag_text if rag_text else "(제공 없음)"}

[출력 JSON 스키마]
{{
  "structure": {{
    "slides_outline": [{{"idx": 1, "title": "...", "key_points": ["...", "..."]}}, ...],
    "logic_flow": {{"good": ["..."], "issues": ["..."]}},
    "missing": ["..."]
  }},
  "content_evaluation": {{
    "scores": {{"기준명": 0}},
    "weighted_total": 0,
    "weak_points": ["..."],
    "comments": ["..."]
  }},
  "evidence_checks": [
    {{"slide_idx": 1, "claim": "...", "evidence": "RAG 요약", "verdict": "일치/불일치/불충분", "note": "..."}}
  ],
  "writing_check": {{
    "typos": [],
    "awkward": [],
    "terminology": [],
    "concise_rewrites": []
  }},
  "feedback": ["핵심 개선 포인트 3~6개"],
  "summary": "한 단락 총평"
}}
주의: 반드시 위 JSON 스키마를 지키고, 숫자는 0~100의 정수만 사용.
"""
    raw = call_openai(
        [{"role": "system", "content": sys}, {"role": "user", "content": usr}],
        max_tokens=1700, temperature=0.2, force_json=True
    )
    result = _parse_json_safely(raw)

    # ---- 5) 맞춤법/문장 체크 보강 1차 ----
    result = _ensure_writing_check(slides_for_llm, result)   # 비면 2차 호출로 채움
    result = _normalize_writing_check(result)                 # 문자열/빈칸 정규화

    # ---- 5-1) 교정안이 없을 때 강제 생성 (run_writing_check) ----
    # before → after/suggest가 비어 있으면 별도 교정 프롬프트로 수정안을 강제 생산
    try:
        wc = result.get("writing_check") or {}
        no_fixes = (
            not any(x.get("after") for x in (wc.get("typos") or []) if isinstance(x, dict)) and
            not any(x.get("suggest") for x in (wc.get("awkward") or []) if isinstance(x, dict)) and
            not any(x.get("line_after") for x in (wc.get("concise_rewrites") or []) if isinstance(x, dict))
        )
        if no_fixes:
            # 슬라이드 텍스트 합치기
            slides_text = "\n".join(
                [f"[{s['idx']}] {s.get('title','')}: {s.get('content','')}" for s in slides_for_llm]
            )
            wc2 = run_writing_check(slides_text)  # ← 아래 함수(교정전용) 사용
            # merge: wc 비면 wc2로 대체
            if wc2:
                result["writing_check"] = wc2
                result = _normalize_writing_check(result)  # 다시 정규화
    except Exception:
        pass

    # ---- 6) 가중 총점 보정 ----
    ce = result.get("content_evaluation", {}) or {}
    scores = (ce.get("scores") or {})
    if not ce.get("weighted_total"):
        weighted = 0
        for name, w in weight_map.items():
            weighted += int(scores.get(name, 0)) * w
        ce["weighted_total"] = round(weighted / denom)
        result["content_evaluation"] = ce

    # ---- 7) RAG 레퍼런스 포함 ----
    if use_rag:
        result["rag_references"] = rag_map
    
    result["slides_for_llm"] = slides_for_llm

    feedback_top = result.get("feedback") or ce.get("weak_points") or []
    return scores, feedback_top, result



def generate_eval_report_docx(uploaded_ppt, criteria, rich):
    """
    구조분석, 내용평가(가중치), 근거검증(RAG), 맞춤법/문장체크, 총평 포함 DOCX.
    """
    try:
        from docx import Document
    except ImportError:
        raise RuntimeError("python-docx가 필요합니다. `pip install python-docx`")

    doc = Document()
    doc.add_heading("PPT 평가 리포트 — LLM Judge + RAG", 0)

    # 총점
    ce = rich.get("content_evaluation", {}) or {}
    total = ce.get("weighted_total", None)
    p = doc.add_paragraph()
    p.add_run("가중 총점: ").bold = True
    p.add_run(str(total) if total is not None else "N/A")

    # 평가기준/가중치
    if criteria:
        doc.add_heading("평가기준 및 가중치", level=1)
        tbl = doc.add_table(rows=1, cols=3)
        hdr = tbl.rows[0].cells
        hdr[0].text = "기준"
        hdr[1].text = "가중치(%)"
        hdr[2].text = "루브릭"
        for c in criteria:
            row = tbl.add_row().cells
            row[0].text = str(c.get("name",""))
            row[1].text = str(c.get("weight",""))
            row[2].text = str(c.get("rubric",""))

    # 내용 평가
    doc.add_heading("내용 평가", level=1)
    scores = (ce.get("scores") or {})
    if scores:
        t2 = doc.add_table(rows=1, cols=2)
        h2 = t2.rows[0].cells
        h2[0].text = "기준"
        h2[1].text = "점수(0~100)"
        for k, v in scores.items():
            r = t2.add_row().cells
            r[0].text = str(k)
            r[1].text = str(v)

    weak = ce.get("weak_points") or []
    if weak:
        doc.add_paragraph("■ 약한 부분")
        for w in weak:
            doc.add_paragraph(w, style="List Bullet")

    comments = ce.get("comments") or []
    if comments:
        doc.add_paragraph("■ 코멘트")
        for c in comments:
            doc.add_paragraph(c, style="List Bullet")

    # 구조 분석
    struct = rich.get("structure", {}) or {}
    doc.add_heading("구조 분석", level=1)
    so = struct.get("slides_outline") or []
    if so:
        doc.add_paragraph("■ 슬라이드별 핵심 정리")
        for s in so:
            idx = s.get("idx","?")
            title = s.get("title","")
            doc.add_paragraph(f"{idx}. {title}", style="List Number")
            for kp in (s.get("key_points") or [])[:8]:
                doc.add_paragraph(kp, style="List Bullet")
    lf = struct.get("logic_flow") or {}
    if lf.get("issues"):
        doc.add_paragraph("■ 논리 흐름 이슈")
        for i in lf["issues"]:
            doc.add_paragraph(i, style="List Bullet")
    if struct.get("missing"):
        doc.add_paragraph("■ 빠진 내용")
        for m in struct["missing"]:
            doc.add_paragraph(m, style="List Bullet")

    # RAG 근거 검증
    if rich.get("evidence_checks") or rich.get("rag_references"):
        doc.add_heading("근거 검증 (RAG)", level=1)

    ev = rich.get("evidence_checks") or []
    if ev:
        doc.add_paragraph("■ 주장 vs 참조 근거")
        for e in ev:
            doc.add_paragraph(
                f"(Slide {e.get('slide_idx','?')}) 주장: {e.get('claim','')}", style="List Bullet"
            )
            doc.add_paragraph(f"  - 근거: {e.get('evidence','')}")
            doc.add_paragraph(f"  - 판단: {e.get('verdict','')}, 비고: {e.get('note','')}")

    rag_map = rich.get("rag_references") or {}
    if rag_map:
        doc.add_paragraph("■ 참고 문헌 목록")
        for sid, refs in rag_map.items():
            doc.add_paragraph(f"- Slide {sid}")
            for r in refs:
                doc.add_paragraph(f"  · {r.get('source','')}: { (r.get('snippet','') or '')[:200] }")

    # 맞춤법/문장 체크 (있으면)
    wc = rich.get("writing_check", {}) or {}
    if any([wc.get("typos"), wc.get("awkward"), wc.get("terminology"), wc.get("concise_rewrites")]):
        doc.add_heading("맞춤법 · 문장 체크", level=1)
        for t in wc.get("typos", []):
            doc.add_paragraph("■ 오타/띄어쓰기")
            doc.add_paragraph(f"{t.get('before','')} → {t.get('after','')} ({t.get('why','')})", style="List Bullet")
        for a in wc.get("awkward", []):
            doc.add_paragraph("■ 어색한 문장 → 제안")
            doc.add_paragraph(f"{a.get('before','')} → {a.get('suggest','')} ({a.get('reason','')})", style="List Bullet")
        for tm in wc.get("terminology", []):
            doc.add_paragraph("■ 용어 일관성")
            doc.add_paragraph(f"{tm.get('term','')}: {tm.get('note','')}", style="List Bullet")
        for rw in wc.get("concise_rewrites", []):
            doc.add_paragraph("■ PPT답게 간결화")
            doc.add_paragraph(
                f"(Slide {rw.get('slide_idx','?')}) {rw.get('line_before','')} → {rw.get('line_after','')}",
                style="List Bullet"
            )

    # 총평
    if rich.get("summary"):
        doc.add_heading("총평", level=1)
        doc.add_paragraph(rich["summary"])

    bio = BytesIO()
    doc.save(bio)
    bio.seek(0)
    return bio

def _shape_texts(shp):
    """개별 shape에서 텍스트 추출 (그룹/테이블 포함)"""
    out = []
    try:
        # 일반 텍스트
        if getattr(shp, "has_text_frame", False):
            txt = shp.text.strip()
            if txt:
                out.append(txt)
        # 테이블
        if getattr(shp, "has_table", False):
            tbl = shp.table
            for r in tbl.rows:
                for c in r.cells:
                    t = c.text.strip()
                    if t:
                        out.append(t)
        # 그룹(내부 도형 순회)
        if shp.shape_type == 6 and hasattr(shp, "shapes"):  # MSO_SHAPE_TYPE.GROUP (6)
            for inner in shp.shapes:
                out.extend(_shape_texts(inner))
    except Exception:
        pass
    return out

def _extract_slide(slide, idx):
    """슬라이드 단위 텍스트 추출(제목 포함)"""
    title = ""
    try:
        if getattr(slide.shapes, "title", None):
            title = (slide.shapes.title.text or "").strip()
    except Exception:
        title = ""

    body = []
    for shp in slide.shapes:
        # 제목은 중복 방지
        if getattr(slide.shapes, "title", None) and shp == slide.shapes.title:
            continue
        body.extend(_shape_texts(shp))

    return {
        "idx": idx,
        "title": title if title else f"슬라이드 {idx}",
        "content": " ".join([t for t in body if t]).strip()
    }

def _collect_text_for_copyedit(slides, limit_chars=3500):
    """맞춤법/문장 교정을 위한 슬라이드 텍스트 모음 (길면 자름)"""
    buf = []
    for s in slides:
        line = f"[{s.get('idx','?')}] {s.get('title','')}: {s.get('content','')}"
        buf.append(line)
    text = "\n".join(buf)
    return text[:limit_chars]

def _heuristic_typos(text):
    """간단 휴리스틱: '샤'같이 드문 음절, 흔한 오타 후보 탐지"""
    suspects = []
    if "제안샤" in text:
        suspects.append({"before": "제안샤", "after": "제안서", "why": "받침/모음 오타"})
    # 필요시 추가 규칙 여기에
    return suspects

def _ensure_writing_check(slides, result):
    """1차 응답에 writing_check가 없으면 2차 교정 호출 + 휴리스틱 보강"""
    wc = result.get("writing_check")
    has_any = isinstance(wc, dict) and any(wc.get(k) for k in ("typos","awkward","terminology","concise_rewrites"))
    if has_any:
        return result

    text = _collect_text_for_copyedit(slides)
    sys = "You are a professional Korean copy editor. 반드시 JSON만 출력."
    usr = f"""
아래 PPT 텍스트를 교정하라.
- 맞춤법/띄어쓰기: 잘못된 부분과 올바른 수정안을 반드시 모두 제시
- 어색한 문장: 'before'와 'suggest'(자연스러운 문장)를 반드시 제시
- 용어 일관성: 혼용된 용어는 대표 표현과 설명 제시
- 간결화: PPT답게 짧고 임팩트있게 수정 제안
각 항목은 'before'(원문)과 'after' 또는 'suggest'(수정안)를 반드시 포함해야 한다.
수정안이 없으면 빈칸 두지 말고 원문을 그대로 넣는다.
반드시 JSON만 출력해라.
예시 오탈자: "제안샤" → "제안서" (받침/모음 오타)

[텍스트]
{text}

출력 JSON 형식:
{{
  "typos": [{{"before":"...","after":"...","why":"..."}}],
  "awkward": [{{"before":"...","suggest":"...","reason":"..."}}],
  "terminology": [{{"term":"...","note":"...","example_before":"...","example_after":"..."}}],
  "concise_rewrites": [{{"slide_idx":1,"line_before":"...","line_after":"..."}}]
}}
주의:
- 실제 텍스트에서 발견된 사례만 내라.
- '제안샤' 같은 오탈자가 있으면 반드시 포함하라.
"""

    raw2 = call_openai(
        [{"role":"system","content":sys}, {"role":"user","content":usr}],
        max_tokens=1200, temperature=0.0, force_json=True
    )
    try:
        wc2 = _parse_json_safely(raw2)
    except Exception:
        wc2 = {"typos": [], "awkward": [], "terminology": [], "concise_rewrites": []}

    # 휴리스틱으로 누락 보강
    h = _heuristic_typos(text)
    if h:
        typos = wc2.get("typos") or []
        # 중복 방지
        existing = { (t.get("before"), t.get("after")) for t in typos }
        for t in h:
            key = (t["before"], t["after"])
            if key not in existing:
                typos.append(t)
        wc2["typos"] = typos

    # 필수 키 보장
    for k in ("typos","awkward","terminology","concise_rewrites"):
        if k not in wc2 or not isinstance(wc2.get(k), list):
            wc2[k] = []

    result["writing_check"] = wc2
    return result

def _norm_wc_item(item, kind):
    # kind: "typos"|"awkward"|"terminology"|"concise"
    if isinstance(item, dict):
        return item
    if isinstance(item, str):
        if kind == "typos":
            return {"before": item, "after": "", "why": ""}
        if kind == "awkward":
            return {"before": item, "suggest": "", "reason": ""}
        if kind == "terminology":
            return {"term": item, "note": ""}
        # concise
        return {"slide_idx": "", "line_before": item, "line_after": ""}
    # unknown type → 문자열화
    s = str(item)
    if kind == "typos":
        return {"before": s, "after": "", "why": ""}
    if kind == "awkward":
        return {"before": s, "suggest": "", "reason": ""}
    if kind == "terminology":
        return {"term": s, "note": ""}
    return {"slide_idx": "", "line_before": s, "line_after": ""}

def _normalize_writing_check(result):
    wc = result.get("writing_check") or {}
    typos = wc.get("typos") or []
    awkward = wc.get("awkward") or []
    terms = wc.get("terminology") or []
    concise = wc.get("concise_rewrites") or []

    wc["typos"] = [_norm_wc_item(x, "typos") for x in typos]
    wc["awkward"] = [_norm_wc_item(x, "awkward") for x in awkward]
    wc["terminology"] = [_norm_wc_item(x, "terminology") for x in terms]
    wc["concise_rewrites"] = [_norm_wc_item(x, "concise") for x in concise]

    result["writing_check"] = wc
    return result

def run_writing_check(slides_text: str):
    """
    슬라이드 텍스트를 받아 맞춤법/띄어쓰기/어색한 문장/용어/간결화 교정 JSON을 생성.
    반드시 before→after/suggest를 채우도록 프롬프트 강화.
    """
    sys = "You are a professional Korean copy editor. 반드시 JSON만 출력한다."
    usr = f"""
아래 PPT 텍스트에서 맞춤법/띄어쓰기 오류, 어색한 문장, 용어 일관성, 간결화를 교정하라.
각 항목은 'before'(원문)과 'after' 또는 'suggest'(수정안)를 반드시 포함해야 한다.
수정안이 없으면 빈칸 두지 말고 원문을 그대로 넣는다.

출력 JSON 형식:
{{
  "typos": [{{"before":"...","after":"...","why":"..."}}]],
  "awkward": [{{"before":"...","suggest":"...","reason":"..."}}]],
  "terminology": [{{"term":"...","note":"...","example_before":"...","example_after":"..."}}]],
  "concise_rewrites": [{{"slide_idx":1,"line_before":"...","line_after":"..."}}]]
}}

[검사할 텍스트]
{slides_text}
"""
    raw = call_openai(
        [{"role": "system", "content": sys}, {"role": "user", "content": usr}],
        max_tokens=1200, temperature=0.0, force_json=True
    )
    out = _parse_json_safely(raw)

    # 필수 배열 보장
    for key in ["typos", "awkward", "terminology", "concise_rewrites"]:
        if key not in out or not isinstance(out[key], list):
            out[key] = []
    return out

_B64_LIKE = re.compile(r"^[A-Za-z0-9+/_-]{6,}=?$")

def _pick_display_source(v):
    """
    Search 문서에서 보기 좋은 display name 뽑기:
    title → metadata/title → filename → id (단, id가 base64/난수면 '문서 ###' 형식으로 대체)
    """
    cand = (
        v.get("title")
        or (v.get("metadata") or {}).get("title")
        or v.get("fileName")
        or v.get("filename")
        or v.get("id")
        or "(제목 없음)"
    )
    # id나 난수처럼 보이면 문서 ### 로 매핑
    if _B64_LIKE.match(str(cand)) or str(cand).lower().startswith(("id:", "doc:", "aXRw")):
        return None  # 호출부에서 순번으로 "문서 001" 붙이게 함
    return cand

def _pick_url(v):
    # 색인할 때 원본 URL/경로를 metadata에 넣어둔 경우 표시
    return (
        v.get("url")
        or (v.get("metadata") or {}).get("url")
        or (v.get("metadata") or {}).get("source")
        or ""
    )

def _pick_snippet(v):
    # semantic captions 우선, 없으면 content
    caps = v.get("@search.captions")
    if isinstance(caps, list) and caps:
        t = caps[0].get("text") or ""
        if t:
            return t
    return v.get("content") or v.get("summary") or ""

def render_rag_for_slide(st, slide_idx: int, title: str, refs: list, max_len=140):
    """
    Streamlit에서 슬라이드 한 개의 RAG 결과를 보기 좋게 렌더링
    refs: [{display, url, snippet, raw}, ...]
    """
    if not refs:
        return

    st.markdown(f"#### Slide {slide_idx} · *{title or ''}*")
    for j, r in enumerate(refs, start=1):
        disp = r.get("display", f"문서 {j:03d}")
        url = r.get("url") or ""
        snip = r.get("snippet","").strip()
        if len(snip) > max_len:
            snip = snip[:max_len] + " …"

        # 제목 줄 (링크가 있으면 링크)
        if url:
            st.markdown(f"- **[{disp}]({url})**  ")
        else:
            st.markdown(f"- **{disp}**  ")

        # 스니펫 (인용 스타일)
        st.caption(f"“{snip}”")

def render_rag_section(st, slides_for_llm: list, rag_map: dict):
    """
    슬라이드 전체에 대해 RAG 섹션 렌더링.
    slides_for_llm: [{idx,title,content}, ...]
    rag_map: { idx: [refs...] }
    """
    for s in slides_for_llm:
        sid = s["idx"]
        refs = rag_map.get(sid, [])
        if refs:
            render_rag_for_slide(st, sid, s.get("title",""), refs)

#심사 시뮬레이션 수행
# 30개+ 역할 풀
def get_judge_roles():
    return [
        # 투자
        "VC 파트너","엔젤 투자자","PE 매니저","액셀러레이터 심사역","코퍼레이트 벤처캐피탈(CVC)",
        # 기업 임원
        "CEO","CTO","CMO","CFO","COO","사업개발 이사","프로덕트 총괄",
        # 정부/공공
        "과제 평가위원","정책 담당자","연구소장","공공기관 심사위원",
        # 전문가
        "전략 컨설턴트","법률 전문가(변호사)","회계사","데이터 사이언티스트","AI 리서처","클라우드 아키텍트",
        # 기타
        "고객 대표","파트너사 매니저","엑스퍼트 멘토","세일즈 디렉터","마케팅 디렉터",
        "보안 아키텍트","UX 디자이너","DevOps 리드","HR 리더","교육/EdTech 전문가",
    ]

def default_specialties():
    return {
        "기술": ["AI/ML","블록체인","IoT","바이오","핀테크","게임"],
        "산업": ["B2B SaaS","이커머스","헬스케어","교육","물류"],
        "기능": ["마케팅","영업","개발","운영","재무","HR"],
        "지역": ["국내","글로벌","미국","일본","동남아","유럽"]
    }

def build_judge_system_prompt(profile: dict, global_rubric_text: str):
    """
    profile: {
      id, name, role, style_carefulness, style_question, style_focus, style_tone,
      specialties (list), exp_years, company_size, invest_stage, success_ratio, persona_text,
      top_priority, second_priority, ignore_items
    }
    """
    role = profile.get("role","심사위원")
    style_care = profile.get("style_carefulness","보통")
    style_q = profile.get("style_question","논리적")
    style_focus = profile.get("style_focus","큰 그림 중시")
    style_tone = profile.get("style_tone","직설적")
    specs = ", ".join(profile.get("specialties") or [])
    exp = profile.get("career_years","N/A")
    size = profile.get("company_size","N/A")
    persona = profile.get("persona_text","").strip()

    pri1 = ", ".join(profile.get("priority_1", []) or [])
    pri2 = ", ".join(profile.get("priority_2", []) or [])
    pri3 = ", ".join(profile.get("priority_3", []) or [])
    pri4 = ", ".join(profile.get("priority_4", []) or [])

    return f"""너는 한국어로 질의하는 **{role}** 심사위원이다. 반드시 JSON 형식으로 출력하라.
- 성격/평가 스타일: 까다로움({style_care}), 질문 스타일({style_q}), 관심 영역({style_focus}), 의사소통({style_tone})
- 전문 분야: {specs}
- 경력: {exp}, 회사 규모: {size}
- 페르소나 노트: {persona or "(없음)"}

[평가 우선순위]
- 1순위: {pri1 or "(미지정)"} 
- 2순위: {pri2 or "(미지정)"}
- 3순위: {pri3 or "(미지정)"} 
- 4순위: {pri4 or "(미지정)"}

[평가기준(요약/루브릭)]
{global_rubric_text or "(입력 없음. 일반 심사 기준 적용)"}

행동 지침:
- 너는 같은 인물로서 대화 내내 일관된 관점과 어투를 유지한다.
- 초기에는 제안서 개요 확인 → 심층 질문(항목별/추궁) → 종합 평가/마무리 순서로 진행한다.
- 답변이 모호하면 구체화를 요구한다. 우수하면 칭찬 후 다음 단계로 넘어간다.
- 항상 1~2문단 이내로 간결하게 질문/피드백한다.
- JSON 객체로만 응답한다.
"""

def judge_greet_and_first_impression(profile: dict, outline_text: str) -> str:
    """
    첫 메시지: 자연스러운 자기소개 + 제안서 첫인상 한두 줄 + 첫 질문 '한 개'만.
    절대 여러 질문을 한 번에 하지 말고, 질문은 반드시 하나로 마무리.
    """
    role = profile.get("role") or "심사위원"
    name = profile.get("name") or ""
    strict = profile.get("style_strict", "보통")
    comm = profile.get("style_comm", "직설적")
    specs = ", ".join(profile.get("specialties", [])[:3]) or "일반"
    years = profile.get("career_years", "")
    company = profile.get("company_size", "")

    sys = (
        "You are a Korean pitch judge. Speak naturally in short sentences. "
        "NEVER output JSON. Keep one question per turn. Keep it concise."
    )
    usr = f"""
다음은 발표 개요(요약)입니다. 이를 바탕으로 채팅형으로만 응답하세요.

[심사위원 프로필]
- 역할: {role} {f"({name})" if name else ""}
- 성향: 까다로움 {strict}, 의사소통 {comm}
- 전문분야: {specs}
- 경력/회사: {years}, {company}

[발표 개요]
{outline_text}

규칙:
1) 첫 메시지에는 '자연스러운 자기소개 한 줄' + '첫인상 한두 줄' + '질문 1개만'을 말한다.
2) 질문은 반드시 하나. 다음 턴에서 사용자의 답을 기다린다.
3) 존대말, 간결한 문장.
"""
    return aoai_chat_text(
        [{"role":"system","content":sys},{"role":"user","content":usr}],
        max_tokens=280, temperature=0.4
    )


def judge_next_turn(profile: dict, outline_text: str, history: list[dict]) -> str:
    """
    사용자의 직전 답변을 포함한 대화 내역(history)을 보고
    '한 개'의 다음 질문만 이어서 묻는다. 필요하면 짧은 피드백 한 줄 + 질문 1개.
    history: [{"role":"assistant"|"user","content":...}, ...]
    """
    role = profile.get("role") or "심사위원"
    strict = profile.get("style_strict", "보통")
    comm = profile.get("style_comm", "직설적")
    specs = ", ".join(profile.get("specialties", [])[:3]) or "일반"

    # 최근 몇 턴만 잘라 주입 (LLM 안정성)
    tail = history[-8:] if len(history) > 8 else history
    chat_dump = "\n".join([f"{m['role'].upper()}: {m['content']}" for m in tail])

    sys = (
        "You are a Korean pitch judge continuing a Q&A. "
        "NEVER output JSON. One short question per turn. If the answer was vague, ask to quantify."
    )
    usr = f"""
[심사위원 프로필]
- 역할: {role}, 성향: 까다로움 {strict}, 의사소통 {comm}, 전문분야: {specs}

[발표 개요]
{outline_text}

[최근 대화]
{chat_dump}

규칙:
1) 사용자의 직전 답을 반영해 짧은 코멘트 한 줄(선택) + '다음 질문 1개'만 한다.
2) 수치/근거/비교를 요구하는 구체 질문을 우선한다.
3) 너무 길게 쓰지 말고, 자연스러운 채팅 톤.
"""
    return aoai_chat_text(
        [{"role":"system","content":sys},{"role":"user","content":usr}],
        max_tokens=240, temperature=0.5
    )

def make_default_profile():
    return {
        "id": str(uuid.uuid4())[:8],
        "name": "",
        "role": "VC 파트너",
        "style_carefulness": "보통",
        "style_question": "논리적",
        "style_focus": "큰 그림 중시",
        "style_tone": "직설적",
        "specialties": [],
        "exp_years": "10년 이상",
        "company_size": "스타트업",
        "persona_text": "",
        "priority_1": [],
        "priority_2": [],
        "priority_3": [],
        "priority_4": [],
    }

def build_brief_from_slides(slides_for_llm: list, max_bullets: int = 3, limit_chars: int = 1800) -> str:
    """
    slides_for_llm: [{"idx":1,"title":"...","content":"..."}...]
    - title와 content에서 앞부분만 뽑아 발표용 개요 텍스트 생성
    """
    if not slides_for_llm:
        return ""
    lines = []
    for s in slides_for_llm:
        title = (s.get("title") or f"슬라이드 {s.get('idx')}").strip()
        content = (s.get("content") or "").strip().splitlines()
        bullets = []
        # content를 문장 단위로 잘라 2~3개만
        for ch in content:
            ch = ch.strip(" -•·\t")
            if ch:
                bullets.append(ch)
            if len(bullets) >= max_bullets:
                break
        lines.append(f"{s.get('idx')}. {title}")
        for b in bullets:
            lines.append(f"   - {b}")
    brief = "\n".join(lines)
    return brief[:limit_chars]

# utils.py 최상단 근처 (call_openai 아래)에 추가
def aoai_chat_text(messages, max_tokens=400, temperature=0.4):
    """
    평문(자연스러운 채팅)으로 답을 받는다. JSON 강제 금지.
    """
    return call_openai(messages, max_tokens=max_tokens, temperature=temperature, force_json=False)

# === utils.py에 추가 ===
def judge_score_answer(profile: dict, outline_text: str, criteria: list, history: list[dict], user_answer: str):
    """
    마지막 사용자 답변을 기준으로 LLM이 기준별 점수(0~100)와 코멘트를 JSON으로 반환.
    criteria: [{"name": "...", "weight": 30, "rubric": "..."}]
    history: [{"role":"assistant"|"user","content":...}, ...]  # 최근 몇 턴만 보내 안정화
    """
    import json

    crit_lines = []
    for c in (criteria or []):
        nm = (c.get("name") or "").strip()
        wt = int(c.get("weight", 0))
        rb = (c.get("rubric") or "").strip()
        crit_lines.append(f"- {nm}({wt}%): {rb}")
    crit_text = "\n".join(crit_lines) if crit_lines else "- 전반 평가(100%): 논리·명확성"

    tail = history[-8:] if len(history) > 8 else (history or [])
    chat_dump = "\n".join([f"{m['role'].upper()}: {m['content']}" for m in tail])

    sys = (
        "You are a rigorous Korean pitch judge. "
        "Return JSON only. Do not include any extra text or code fences. "
        "The word JSON is explicitly present here to satisfy response_format rules."
    )
    usr = f"""
아래 발표 개요와 대화 흐름, 그리고 사용자의 '직전 답변'을 기준으로 평가하세요.
반드시 'JSON 객체'만 출력합니다. (설명 금지)

[평가기준과 가중치]
{crit_text}

[발표 개요(요약)]
{outline_text}

[최근 대화]
{chat_dump}

[직전 사용자 답변]
{user_answer}

출력 스키마(JSON):
{{
  "scores": {{"기준명": 0}},        // 각 기준 0~100
  "weighted_total": 0,              // 가중치 적용 총점(정수)
  "comments": ["...","..."]         // 짧은 개선 코멘트 2~5개
}}
"""
    raw = call_openai(
        [{"role": "system", "content": sys}, {"role": "user", "content": usr}],
        max_tokens=350, temperature=0.2, force_json=True
    )
    data = _parse_json_safely(raw)

    # weighted_total 보정 (없거나 잘못되면 재계산)
    try:
        weight_sum = sum(int(c.get("weight", 0)) for c in (criteria or [])) or 1
        wt_total = 0
        for c in (criteria or []):
            nm = (c.get("name") or "").strip()
            w = int(c.get("weight", 0))
            s = int((data.get("scores") or {}).get(nm, 0))
            wt_total += s * w
        data["weighted_total"] = int(round(wt_total / weight_sum))
    except Exception:
        data.setdefault("weighted_total", 0)

    data.setdefault("scores", {})
    data.setdefault("comments", [])
    return data
